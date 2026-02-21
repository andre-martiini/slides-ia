import os
import io
import streamlit as st
from typing import List, Optional
from pydantic import BaseModel
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches

# --- CONFIGURAÇÃO ---
API_KEY = "XXX" # Substitua pela sua chave
MODELO = 'gemini-2.5-flash'

class Slide(BaseModel):
    numero: int
    layout: str
    titulo: str
    topicos: List[str]
    prompt_imagem: Optional[str]

class Apresentacao(BaseModel):
    slides: List[Slide]

# --- FUNÇÕES ---
def gerar_conteudo_llm(rascunho: str, qtd_slides: int) -> Apresentacao:
    client = genai.Client(api_key=API_KEY)
    
    # Instrução atualizada para forçar a quantidade de slides
    system_instruction = f"""
    Atue como Especialista em Design de Apresentações.
    Refine o texto para EXATAMENTE {qtd_slides} slides.
    Regras:
    - layouts: 'capa', 'titulo_e_conteudo', 'somente_titulo'.
    - topicos: frases curtas, máx 4 por slide.
    - prompt_imagem: Inglês, focado no tema.
    """

    response = client.models.generate_content(
        model=MODELO,
        contents=rascunho,
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            response_mime_type="application/json",
            response_schema=Apresentacao,
            temperature=0.2
        ),
    )
    return response.parsed

def criar_pptx(apresentacao: Apresentacao) -> io.BytesIO:
    prs = Presentation()
    
    # Configuração para 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts_map = {"capa": 0, "titulo_e_conteudo": 1, "somente_titulo": 5}

    for slide_data in apresentacao.slides:
        idx = layouts_map.get(slide_data.layout, 1)
        slide = prs.slides.add_slide(prs.slide_layouts[idx])

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.titulo

        if slide_data.layout == "titulo_e_conteudo":
            try:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide_data.topicos[0] if slide_data.topicos else ""
                for i in range(1, len(slide_data.topicos)):
                    p = tf.add_paragraph()
                    p.text = slide_data.topicos[i]
            except Exception:
                pass

        if slide_data.prompt_imagem:
            slide.notes_slide.notes_text_frame.text = f"PROMPT DE IMAGEM: {slide_data.prompt_imagem}"

    # Salva em memória (BytesIO) para o Streamlit processar o download
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

# --- INTERFACE STREAMLIT ---
st.set_page_config(page_title="Gerador de PPTX IA", layout="centered")

st.title("Gerador de Apresentações (16:9)")

rascunho = st.text_area("Texto Bruto", height=200, placeholder="Cole seu texto base aqui...")
qtd_slides = st.number_input("Quantidade de Slides", min_value=1, max_value=50, value=5)

if st.button("Gerar Apresentação", type="primary"):
    if not rascunho.strip():
        st.warning("Insira o texto bruto antes de gerar.")
    else:
        with st.spinner("Processando com IA e montando slides..."):
            try:
                dados_slides = gerar_conteudo_llm(rascunho, qtd_slides)
                arquivo_pptx = criar_pptx(dados_slides)
                
                st.success("Apresentação gerada com sucesso!")
                st.download_button(
                    label="📥 Baixar PowerPoint (.pptx)",
                    data=arquivo_pptx,
                    file_name="apresentacao_designer.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"Erro durante a geração: {e}")
