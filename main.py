import os
from typing import List, Optional
from pydantic import BaseModel
from google import genai
from google.genai import types
from pptx import Presentation

# Schema de Validação
class Slide(BaseModel):
    numero: int
    layout: str
    titulo: str
    topicos: List[str]
    prompt_imagem: Optional[str]

class Apresentacao(BaseModel):
    slides: List[Slide]

def gerar_conteudo_llm(rascunho: str, api_key: str) -> Apresentacao:
    client = genai.Client(api_key=api_key)
    
    # Modelo 2.0 Flash: rápido e confiável
    modelo = 'gemini-2.0-flash' 

    system_instruction = """
    Atue como Especialista em Design de Apresentações.
    Refine o texto para slides de impacto.
    Regras:
    - layouts: 'capa', 'titulo_e_conteudo', 'somente_titulo'.
    - topicos: frases curtas, máx 4 por slide.
    - prompt_imagem: inglês, estilo fotorealista ou 3D render, focado no tema do slide.
    """

    response = client.models.generate_content(
        model=modelo,
        contents=rascunho,
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            response_mime_type="application/json",
            response_schema=Apresentacao,
        ),
    )
    return response.parsed

def criar_pptx(apresentacao: Apresentacao, nome_arquivo: str = "apresentacao_vitoria.pptx"):
    prs = Presentation()
    layouts = {"capa": 0, "titulo_e_conteudo": 1, "somente_titulo": 5}

    for slide_data in apresentacao.slides:
        idx = layouts.get(slide_data.layout, 1)
        slide = prs.slides.add_slide(prs.slide_layouts[idx])

        if slide.shapes.title:
            slide.shapes.title.text = slide_data.titulo

        # Tenta preencher o corpo de texto se existir
        try:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_data.topicos[0] if slide_data.topicos else ""
            for i in range(1, len(slide_data.topicos)):
                p = tf.add_paragraph()
                p.text = slide_data.topicos[i]
        except:
            pass

        # Notas do orador com o Prompt
        if slide_data.prompt_imagem:
            slide.notes_slide.notes_text_frame.text = slide_data.prompt_imagem

    prs.save(nome_arquivo)
    print(f"\n✅ Sucesso! Arquivo '{nome_arquivo}' gerado.")

if __name__ == "__main__":
    # COLE SUA NOVA CHAVE AQUI
    MINHA_NOVA_CHAVE = "AIzaSyCKONuumd3OwGV0k2axTCwDG3HDGpwo_5g"
    
    meu_rascunho = "Crie uma apresentação sobre Governança Pública e IA no IFES."
    
    try:
        dados = gerar_conteudo_llm(meu_rascunho, MINHA_NOVA_CHAVE)
        criar_pptx(dados)
    except Exception as e:
        print(f"\n❌ Erro ainda persiste: {e}")